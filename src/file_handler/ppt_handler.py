from pptx import Presentation
from pptx.shapes.graphfrm import GraphicFrame


class PptHandler:
    def __init__(self, translator):
        self.translator = translator

    def set_font(self, shape, font_name):
        if isinstance(shape, GraphicFrame):
            pass
        else:
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.name = font_name

    def translate_payload(self, slide):
        def _translate_text(obj):
            if obj.text and str(obj.text).strip():
                to_translate = str(obj.text).strip()
                translated = self.translator.translate(to_translate)
                print("\t{}: {}->{}".format(type(obj), obj.text, translated))
                if translated:
                    obj.text = translated

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        _translate_text(run)
            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        _translate_text(cell)
            else:
                print("Unsupported type {}".format(shape))

    def translate(self, path):
        presentation = Presentation(path)
        print(len(presentation.slides))
        for slide, id in zip(presentation.slides, range(len(presentation.slides))):
            print("slide", id + 1)
            self.translate_payload(slide)
            id += 1

        presentation.save(path[:-5] + "-tr" + path[-5:])
