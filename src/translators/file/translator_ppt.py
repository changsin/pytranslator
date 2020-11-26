from pptx import Presentation


class TranslatorPpt:
    def __init__(self, translator):
        self.translator = translator

    def set_font(self, shape, font_name):
        for paragraph in shape.text_frame.paragraphs:
            paragraph.font.name = font_name

    def translate_payload(self, slide):
        slide_str = ""
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                if shape.text and str(shape.text).strip():
                    translated_text = self.translator.translate(str(shape.text).strip())
                    slide_str = slide_str + "\t: {}->{}".format(str(shape.text).strip(), translated_text)
                    if translated_text:
                        shape.text = translated_text
                    # self.set_font(shape, 'Gulim')

                    print(slide_str)

    def translate(self, path):
        presentation = Presentation(path)
        for slide in presentation.slides:
            for shape in slide.shapes:

                if hasattr(shape, 'text'):
                    self.translate_payload(slide)

            presentation.save(path[:-5] + "-tr" + path[-5:])
